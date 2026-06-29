from io import BytesIO

from pptx import Presentation

from app.parsers.base import download_from_url


async def parse_pptx(url: str) -> str:
    content = await download_from_url(url)
    prs = Presentation(BytesIO(content))

    slides = []
    for i, slide in enumerate(prs.slides):
        texts = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                texts.append(shape.text.strip())
        slides.append(f"Slide {i + 1}:\n" + "\n".join(texts))

    return "\n\n".join(slides)
